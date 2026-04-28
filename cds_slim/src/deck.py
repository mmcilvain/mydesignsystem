"""
deck.py — template-driven deck builder.

Workflow:
    deck = Deck()                                 # opens the brand template
    deck.add('presentation_title_slide',
             title='Q3 review', subhead='Revenue + product wins',
             date='OCTOBER 2026')
    deck.add('comcast_standard',
             banner='OPERATING REVIEW', title='Where we landed',
             body='Lorem ipsum...')
    deck.save('out.pptx')

The template has all 94 layouts AND all the Comcast embedded fonts, which is
why we open it instead of building from scratch. By default the 95 sample
slides shipped with the template are stripped on open. Set clear_samples=False
to keep them.

For layouts with repeated regions (e.g. Key Points has 5x subheader/body),
pass lists: subheaders=[...], bodies=[...]. List items are mapped in
top-to-bottom, left-to-right order based on the layout's placeholder y-then-x
positions.
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Any, Iterable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER

from . import brand
from .brand import LayoutSpec


# placeholder type values for titles (we match by enum, not raw int, for safety)
_TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


# default location for the template — assumes repo layout
DEFAULT_TEMPLATE_PATH = (
    Path(__file__).resolve().parent.parent / "assets" / "comcast_template.pptx"
)


# kwargs the caller can pass for each semantic role. The classifier below maps
# placeholders to roles based on the layout's sample text (extracted from the
# original .pptx). This keeps the API stable across layouts that share roles
# but differ in placement.
ROLE_PATTERNS: dict[str, list[str]] = {
    "title":          ["click to edit master title", "impact headline", "presentation title", "headline text", "section", "agenda"],
    "banner":         ["banner text"],
    "subhead":        ["standard subhead", "optional subhead", "subhead"],
    "body":           ["click to edit master text", "optional paragraph", "lorem ipsum", "list items", "call-out text"],
    "subtitle":       ["subhead", "subtitle", "subheader"],
    "date":           ["february 2026", "october 2026", "month year", "date"],
    "confidentiality":["private & confidential"],
    "page_number":    ["‹#›"],
    "stat_value":     ["12.3m", "12.3%", "12.3k"],
    "stat_caption":   ["data caption"],
}


def _classify_role(ph_type: str | None, sample: str, name: str) -> str | None:
    """Best-guess role for a placeholder based on its template metadata."""
    s = (sample or "").lower()
    n = (name or "").lower()
    if ph_type in ("title", "ctrTitle"):
        return "title"
    if "slide number" in n or s == "‹#›":
        return "page_number"
    for role, patterns in ROLE_PATTERNS.items():
        for p in patterns:
            if p in s:
                return role
    return None


def _placeholders_for_role(spec: LayoutSpec, role: str) -> list[brand.Placeholder]:
    """All placeholders in the layout matching the given role, in reading order."""
    matched = [p for p in spec.placeholders
               if _classify_role(p.ph_type, p.sample_text, p.name) == role]
    # reading order: top-to-bottom, then left-to-right
    matched.sort(key=lambda p: (p.y_in or 0, p.x_in or 0))
    return matched


# ---------------------------------------------------------------------------
# Slide deletion — python-pptx doesn't expose this, so we go to the XML.
# Reference: standard recipe from python-pptx maintainers.
# ---------------------------------------------------------------------------

def _delete_all_slides(prs) -> None:
    sldIdLst = prs.slides._sldIdLst   # CT_SlideIdList
    rels = prs.part.rels
    # iterate in reverse so list mutation stays consistent
    for sldId in list(sldIdLst):
        rid = sldId.rId
        # drop the rel and the slide part
        prs.part.drop_rel(rid)
        sldIdLst.remove(sldId)


# ---------------------------------------------------------------------------
# Deck
# ---------------------------------------------------------------------------

class Deck:
    """Wraps a python-pptx Presentation seeded from the Comcast brand template."""

    def __init__(self,
                 template_path: str | Path | None = None,
                 clear_samples: bool = True) -> None:
        path = Path(template_path) if template_path else DEFAULT_TEMPLATE_PATH
        if not path.exists():
            raise FileNotFoundError(
                f"template not found at {path}. "
                "Either pass template_path explicitly or use comcast_design_system.fallback "
                "to build a deck without the .pptx."
            )
        self._prs = Presentation(str(path))
        if clear_samples:
            _delete_all_slides(self._prs)

    @property
    def presentation(self):
        """Access the raw python-pptx Presentation if you need to do something custom."""
        return self._prs

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    # ------------------------------------------------------------------
    # adding slides
    # ------------------------------------------------------------------

    def add(self, layout_id: str | int, **content: Any):
        """
        Add a slide using a layout (by slug, name, or slot index) and fill placeholders.

        Recognised content kwargs (use whichever apply for the layout):
            title, banner, subhead, body, subtitle, date, confidentiality
            subheaders=[...]   for layouts with repeated subheader rows
            bodies=[...]       for layouts with repeated body rows
            stat_values=[...]  for big-stat layouts (Data_3, Data_4, Data_6, Data_8)
            stat_captions=[...]
            placeholders={idx: 'text', ...}  raw idx-based override

        Returns the python-pptx Slide so callers can do anything else they want.
        """
        spec = brand.layout(layout_id)
        layout_obj = self._prs.slide_layouts[spec.slot_index]
        slide = self._prs.slides.add_slide(layout_obj)

        # build idx -> placeholder map for the new slide
        idx_to_ph = {}
        for ph in slide.placeholders:
            idx_to_ph[ph.placeholder_format.idx] = ph

        # raw idx overrides win — apply first so role-based fills don't clobber them
        raw_idx = content.pop("placeholders", {}) or {}
        for idx, text in raw_idx.items():
            if idx in idx_to_ph:
                _set_text(idx_to_ph[idx], str(text))

        # singular roles (title/banner/subhead/body/etc)
        for role in ("title", "banner", "subhead", "subtitle", "date", "body", "confidentiality"):
            if role not in content:
                continue
            value = content.pop(role)
            if role == "title":
                # title is special — python-pptx assigns idx=0 and type=TITLE/CENTER_TITLE.
                # the layout XML often lacks an explicit idx attribute on the title shape.
                title_ph = None
                for ph in slide.placeholders:
                    if ph.placeholder_format.type in _TITLE_TYPES:
                        title_ph = ph
                        break
                if title_ph is None and 0 in idx_to_ph:
                    title_ph = idx_to_ph[0]
                if title_ph is not None:
                    _set_text(title_ph, str(value))
                continue

            phs = _placeholders_for_role(spec, role)
            if not phs:
                continue
            target_idx = phs[0].ph_idx
            if target_idx is not None and int(target_idx) in idx_to_ph:
                _set_text(idx_to_ph[int(target_idx)], str(value))

        # plural roles
        plural_map = {
            "subheaders":   "subtitle",
            "bodies":       "body",
            "stat_values":  "stat_value",
            "stat_captions":"stat_caption",
        }
        for kw, role in plural_map.items():
            if kw not in content:
                continue
            values: list[str] = list(content.pop(kw))
            phs = _placeholders_for_role(spec, role)
            for value, ph_spec in zip(values, phs):
                if ph_spec.ph_idx is None:
                    continue
                idx = int(ph_spec.ph_idx)
                if idx in idx_to_ph:
                    _set_text(idx_to_ph[idx], str(value))

        # anything left in content is an unknown kwarg — surface it loudly
        if content:
            raise TypeError(f"unknown content kwargs for layout {spec.slug!r}: {list(content)}. "
                            f"valid kwargs: title/banner/subhead/body/subtitle/date/"
                            f"subheaders/bodies/stat_values/stat_captions/placeholders")

        return slide

    def add_blank(self):
        """Shortcut: add a true blank slide (uses 'BLANK' layout, slot 89)."""
        return self.add("blank")

    def add_closing(self):
        """Shortcut: closing slide (uses 'CLOSING' slot 90)."""
        return self.add("closing")

    # ------------------------------------------------------------------
    # introspection
    # ------------------------------------------------------------------

    def describe_layout(self, layout_id: str | int) -> str:
        """Print the placeholder schema for a layout — handy when scripting."""
        spec = brand.layout(layout_id)
        lines = [f"\n[{spec.slot_index}] {spec.name}  slug={spec.slug}  family={spec.family}"]
        for p in spec.placeholders:
            role = _classify_role(p.ph_type, p.sample_text, p.name) or "?"
            geom = (f"x={p.x_in:.2f} y={p.y_in:.2f} w={p.w_in:.2f} h={p.h_in:.2f}"
                    if p.x_in is not None else "(inherits master)")
            sample = p.sample_text[:50] + "..." if len(p.sample_text) > 50 else p.sample_text
            lines.append(f"  idx={p.ph_idx or '-':>3}  role={role:14s}  type={p.ph_type or '-':10s}  {geom}")
            if sample:
                lines.append(f"        sample: {sample!r}")
        return "\n".join(lines)

    # ------------------------------------------------------------------
    # save
    # ------------------------------------------------------------------

    def save(self, path: str | Path) -> Path:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(path))
        return path


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------

def _set_text(placeholder, text: str) -> None:
    """Replace placeholder text. Layout and theme inheritance handles font/size/color
    so we don't need to manually preserve runs."""
    if not placeholder.has_text_frame:
        return
    placeholder.text_frame.text = text
