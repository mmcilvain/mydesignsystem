"""
fallback.py — generate decks WITHOUT the .pptx template.

When you can't ship the .pptx (CI runners, restricted envs, etc), use this.
It builds slides from scratch using brand colors + tokens, with the Aptos
fallback fonts. Output won't match the .pptx pixel-for-pixel — no peacock
graphics, no embedded Comcast New Vision — but it'll still be on-palette
and structured.

Scope: the 7 most common patterns. If you need bento, mockups, gantts, etc.
either install the .pptx and use deck.Deck, or extend this module.
"""
from __future__ import annotations

from pathlib import Path
from typing import Sequence

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import brand


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# common positions in inches — copied from brand.json grid for parity
_BANNER = brand.REGION_BANNER
_TITLE = brand.REGION_TITLE
_SUBHEAD = brand.REGION_SUBHEAD
_PAGE_NUM = brand.REGION_PAGE_NUMBER
_CONF = brand.REGION_CONFIDENTIALITY


class BlankDeck:
    """Builds a Comcast-branded 16:9 deck without the .pptx template.

    All shapes are drawn from scratch. Fonts fall back to Aptos / Aptos Display
    when Comcast New Vision isn't available. Colors come from brand.COLORS.
    """

    def __init__(self) -> None:
        self._prs = Presentation()
        # 16:9 widescreen — 13.333 x 7.5
        self._prs.slide_width = Emu(brand.SLIDE_WIDTH_EMU)
        self._prs.slide_height = Emu(brand.SLIDE_HEIGHT_EMU)
        self._counter = 0  # for page numbers

    @property
    def presentation(self):
        return self._prs

    @property
    def slide_count(self) -> int:
        return len(self._prs.slides)

    # ------------------------------------------------------------------
    # title slide — big title, optional subhead + date
    # ------------------------------------------------------------------

    def add_title_slide(self, title: str, subhead: str = "", date: str = "") -> None:
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])  # blank
        self._counter += 1

        # date pill (top-left corner if provided)
        if date:
            self._add_text_box(slide, x=0.39, y=2.65, w=2.84, h=0.42, text=date,
                               font=brand.FONTS.banner, size_pt=11, bold=True,
                               color=brand.COLORS.black, align=PP_ALIGN.CENTER,
                               anchor=MSO_ANCHOR.MIDDLE)
            # rounded rect outline behind it
            self._add_pill_outline(slide, x=0.39, y=2.65, w=2.84, h=0.42)

        # main title — large, lower-left
        self._add_text_box(slide, x=0.49, y=3.37, w=10.0, h=2.61, text=title,
                           font=brand.FONTS.heading, size_pt=60, bold=False,
                           color=brand.COLORS.black, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

        # subhead at the bottom-left
        if subhead:
            self._add_text_box(slide, x=0.49, y=6.09, w=8.0, h=0.48, text=subhead,
                               font=brand.FONTS.body, size_pt=18, bold=False,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT)

        # confidentiality marker
        self._add_text_box(slide, x=10.55, y=7.13, w=2.4, h=_CONF.h,
                           text="PRIVATE & CONFIDENTIAL",
                           font=brand.FONTS.body, size_pt=8, bold=False,
                           color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)

    # ------------------------------------------------------------------
    # section divider — colored band + section title
    # ------------------------------------------------------------------

    def add_section_divider(self, title: str, eyebrow: str = "",
                            accent_hex: str = None) -> None:
        accent = accent_hex or brand.COLORS.yellow
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        self._counter += 1

        # full-bleed accent band on the left third
        self._add_rect(slide, x=0, y=0, w=4.0, h=brand.SLIDE_HEIGHT_IN, fill_hex=accent)

        # eyebrow
        if eyebrow:
            self._add_text_box(slide, x=4.5, y=2.8, w=8.5, h=0.4, text=eyebrow.upper(),
                               font=brand.FONTS.banner, size_pt=12, bold=True,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT)

        # title
        self._add_text_box(slide, x=4.5, y=3.3, w=8.5, h=2.0, text=title,
                           font=brand.FONTS.heading, size_pt=54, bold=False,
                           color=brand.COLORS.black, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # ------------------------------------------------------------------
    # standard content slide — banner / title / subhead / body
    # ------------------------------------------------------------------

    def add_standard_slide(self, title: str, banner: str = "", subhead: str = "",
                           body: str = "") -> None:
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        self._counter += 1
        self._draw_chrome(slide, banner)

        # title
        self._add_text_box(slide, x=_TITLE.x, y=_TITLE.y, w=_TITLE.w, h=_TITLE.h, text=title,
                           font=brand.FONTS.heading, size_pt=28, bold=False,
                           color=brand.COLORS.black, align=PP_ALIGN.LEFT)
        # subhead
        if subhead:
            self._add_text_box(slide, x=_SUBHEAD.x, y=_SUBHEAD.y, w=_SUBHEAD.w, h=_SUBHEAD.h,
                               text=subhead, font=brand.FONTS.body, size_pt=14,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT)
        # body
        if body:
            self._add_text_box(slide, x=0.37, y=2.19, w=12.59, h=4.81, text=body,
                               font=brand.FONTS.body, size_pt=14,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT,
                               anchor=MSO_ANCHOR.TOP)

    # ------------------------------------------------------------------
    # data slide — 1-4 big-stat callouts
    # ------------------------------------------------------------------

    def add_data_slide(self, title: str, stats: Sequence[tuple[str, str]],
                       banner: str = "") -> None:
        if not (1 <= len(stats) <= 4):
            raise ValueError("stats must contain 1-4 (value, caption) tuples")
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        self._counter += 1
        self._draw_chrome(slide, banner)

        self._add_text_box(slide, x=_TITLE.x, y=_TITLE.y, w=_TITLE.w, h=_TITLE.h, text=title,
                           font=brand.FONTS.heading, size_pt=28,
                           color=brand.COLORS.black, align=PP_ALIGN.LEFT)

        # split content area into N equal columns
        n = len(stats)
        margin = 0.5
        gap = 0.2
        total_w = brand.SLIDE_WIDTH_IN - 2 * margin - gap * (n - 1)
        col_w = total_w / n
        y_value = 2.8
        y_caption = 5.4
        for i, (value, caption) in enumerate(stats):
            x = margin + i * (col_w + gap)
            self._add_text_box(slide, x=x, y=y_value, w=col_w, h=2.4, text=value,
                               font=brand.FONTS.heading, size_pt=72, bold=False,
                               color=brand.COLORS.black, align=PP_ALIGN.CENTER,
                               anchor=MSO_ANCHOR.MIDDLE)
            self._add_text_box(slide, x=x, y=y_caption, w=col_w, h=0.8, text=caption,
                               font=brand.FONTS.body, size_pt=14,
                               color=brand.COLORS.black, align=PP_ALIGN.CENTER)

    # ------------------------------------------------------------------
    # key points — 3-6 (header, body) rows on the right, optional title on the left
    # ------------------------------------------------------------------

    def add_key_points_slide(self, title: str, points: Sequence[tuple[str, str]],
                             banner: str = "") -> None:
        if not (1 <= len(points) <= 6):
            raise ValueError("points must contain 1-6 (header, body) tuples")
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        self._counter += 1
        self._draw_chrome(slide, banner)

        # title on the left, points stacked on the right
        self._add_text_box(slide, x=0.38, y=0.89, w=5.5, h=2.7, text=title,
                           font=brand.FONTS.heading, size_pt=36,
                           color=brand.COLORS.black, align=PP_ALIGN.LEFT,
                           anchor=MSO_ANCHOR.TOP)

        n = len(points)
        top = 1.05
        bottom = 7.0
        row_h = (bottom - top) / n
        for i, (header, body) in enumerate(points):
            y = top + i * row_h
            self._add_text_box(slide, x=6.89, y=y, w=5.95, h=0.25, text=header,
                               font=brand.FONTS.heading, size_pt=14, bold=True,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT)
            self._add_text_box(slide, x=6.89, y=y + 0.30, w=5.95, h=row_h - 0.35,
                               text=body, font=brand.FONTS.body, size_pt=12,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT,
                               anchor=MSO_ANCHOR.TOP)

    # ------------------------------------------------------------------
    # closing — minimalist
    # ------------------------------------------------------------------

    def add_closing_slide(self, message: str = "Thank you") -> None:
        slide = self._prs.slides.add_slide(self._prs.slide_layouts[6])
        self._add_text_box(slide, x=0, y=3.0, w=brand.SLIDE_WIDTH_IN, h=1.5,
                           text=message,
                           font=brand.FONTS.heading, size_pt=54,
                           color=brand.COLORS.black, align=PP_ALIGN.CENTER,
                           anchor=MSO_ANCHOR.MIDDLE)

    # ------------------------------------------------------------------
    # save
    # ------------------------------------------------------------------

    def save(self, path: str | Path) -> Path:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self._prs.save(str(path))
        return path

    # ------------------------------------------------------------------
    # internals
    # ------------------------------------------------------------------

    def _draw_chrome(self, slide, banner_text: str) -> None:
        """Banner top-left + page number top-right + confidentiality marker. Same on every content slide."""
        if banner_text:
            self._add_text_box(slide, x=_BANNER.x, y=_BANNER.y, w=_BANNER.w, h=_BANNER.h,
                               text=banner_text.upper(),
                               font=brand.FONTS.banner, size_pt=9, bold=True,
                               color=brand.COLORS.black, align=PP_ALIGN.LEFT)
        # confidentiality
        self._add_text_box(slide, x=10.55, y=_CONF.y, w=2.4, h=_CONF.h,
                           text="PRIVATE & CONFIDENTIAL",
                           font=brand.FONTS.body, size_pt=8,
                           color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.RIGHT)
        # page number
        self._add_text_box(slide, x=_PAGE_NUM.x, y=_PAGE_NUM.y,
                           w=_PAGE_NUM.w, h=_PAGE_NUM.h, text=str(self._counter),
                           font=brand.FONTS.body, size_pt=9,
                           color=brand.COLORS.black, align=PP_ALIGN.RIGHT)

    def _add_text_box(self, slide, x: float, y: float, w: float, h: float, text: str,
                      *, font: str, size_pt: float, color: RGBColor,
                      bold: bool = False, italic: bool = False,
                      align: PP_ALIGN = PP_ALIGN.LEFT,
                      anchor: MSO_ANCHOR = MSO_ANCHOR.TOP) -> None:
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.margin_left = Inches(0)
        tf.margin_right = Inches(0)
        tf.margin_top = Inches(0)
        tf.margin_bottom = Inches(0)
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color if isinstance(color, RGBColor) else _hex_to_rgb(color)

    def _add_rect(self, slide, x: float, y: float, w: float, h: float, fill_hex: str) -> None:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.solid()
        rect.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
        rect.line.fill.background()  # no border

    def _add_pill_outline(self, slide, x: float, y: float, w: float, h: float) -> None:
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(w), Inches(h))
        rect.fill.background()
        rect.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        rect.line.width = Pt(0.75)
